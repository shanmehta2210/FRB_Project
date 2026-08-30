"""Build PPTX: population diagnostics first, then confirmed host panels.

Panel / plot PNG files on disk are **never modified**. For slide embedding
only, RGBA is flattened to RGB in memory and pictures are placed at a fixed
high DPI (exact pixel→inch sizing when they fit; otherwise letterboxed to the
slide while keeping the full pixel blob).

See ``VISUAL_PANELS.md``.
"""

from __future__ import annotations

import io
import json
import re
import sys
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

VER = Path(__file__).resolve().parent
if str(VER) not in sys.path:
    sys.path.insert(0, str(VER))

import vercommon as vc  # noqa: E402

PANELS = Path(vc.OUT_ROOT) / "panels"
PLOTS = Path(vc.OUT_ROOT) / "plots"
HC = VER / "host_confirmation.csv"
POP_SUMMARY = Path(vc.OUT_ROOT) / "tables" / "population_summary.json"
OUT = Path(vc.OUT_ROOT) / "confirmed_fit_panels.pptx"

EMBED_DPI = 200
TITLE_H_IN = 0.55
MARGIN_IN = 0.12
EMU_PER_INCH = 914400

_ALT = re.compile(
    r"outputs/panels/([A-Za-z0-9]+_(?:n1_sky|n1|sky|psf))\.png"
)

# One population plot per slide (do not pack onto a single page).
POPULATION_SLIDES: list[tuple[str, str, str]] = [
    (
        "population_diagnostics.png",
        "Population diagnostics (science cut)",
        "χ² vs SNR · RFF vs Re/FWHM · q vs e_PSF · |ΔPA| · AstroPhot vs GALFIT · isophote vs GALFIT",
    ),
    (
        "mag_leakage.png",
        "Magnitude leakage (cohort OLS)",
        "Δm vs Re, n, sky offset, χ² — strong Δm–Re trend is the main photometric systematic",
    ),
    (
        "dq_comparison.png",
        "Geometry Δq distributions",
        "Fourier δq · sky ±1σ Δq · AstroPhot Δq  (production metrics; not alternate-leg)",
    ),
    (
        "contact_sheet.png",
        "Residual contact sheet (all 64)",
        "Confirmed-leg / PSF-only residuals where gated; production otherwise. ±5σ, crop 3 Re",
    ),
]


def panel_for(frb: str, notes: str) -> Path:
    m = _ALT.search(notes or "")
    if m:
        p = PANELS / f"{m.group(1)}.png"
        if p.is_file():
            return p
    return PANELS / f"{frb}.png"


def _embed_stream(path: Path) -> tuple[io.BytesIO, int, int]:
    """Load PNG for embedding; flatten alpha → white RGB without touching disk."""
    with Image.open(path) as im:
        im.load()
        w, h = im.size
        if im.mode == "RGBA":
            bg = Image.new("RGB", im.size, (255, 255, 255))
            bg.paste(im, mask=im.split()[3])
            rgb = bg
        elif im.mode != "RGB":
            rgb = im.convert("RGB")
        else:
            rgb = im.copy()
    buf = io.BytesIO()
    rgb.save(buf, format="PNG", compress_level=3, optimize=False)
    buf.seek(0)
    return buf, w, h


def _px_to_emu(px: int, dpi: int = EMBED_DPI) -> int:
    return int(round(px / dpi * EMU_PER_INCH))


def _configure_slide_size(prs: Presentation, px_w: int, px_h: int) -> None:
    prs.slide_width = Emu(_px_to_emu(px_w) + int(2 * MARGIN_IN * EMU_PER_INCH))
    prs.slide_height = Emu(
        _px_to_emu(px_h) + int((TITLE_H_IN + 2 * MARGIN_IN) * EMU_PER_INCH)
    )


def _fit_picture_box(
    px_w: int, px_h: int, max_w_emu: int, max_h_emu: int
) -> tuple[int, int]:
    """Letterbox full-res image into max box; prefer native 200 DPI if it fits."""
    native_w = _px_to_emu(px_w)
    native_h = _px_to_emu(px_h)
    if native_w <= max_w_emu and native_h <= max_h_emu:
        return native_w, native_h
    scale = min(max_w_emu / native_w, max_h_emu / native_h)
    return int(native_w * scale), int(native_h * scale)


def add_image_slide(
    prs: Presentation,
    title: str,
    img: Path,
    *,
    title_color: RGBColor | None = None,
    subtitle: str = "",
) -> None:
    stream, px_w, px_h = _embed_stream(img)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw = int(prs.slide_width)
    sh = int(prs.slide_height)

    box = slide.shapes.add_textbox(
        Emu(int(MARGIN_IN * EMU_PER_INCH)),
        Emu(int(0.06 * EMU_PER_INCH)),
        Emu(sw - int(2 * MARGIN_IN * EMU_PER_INCH)),
        Emu(int(TITLE_H_IN * EMU_PER_INCH)),
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.name = "Calibri"
    p.font.color.rgb = title_color or RGBColor(0x1A, 0x1A, 0x1A)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(10)
        p2.font.name = "Calibri"
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

    avail_w = sw - int(2 * MARGIN_IN * EMU_PER_INCH)
    top = int((MARGIN_IN + TITLE_H_IN) * EMU_PER_INCH)
    avail_h = sh - top - int(MARGIN_IN * EMU_PER_INCH)
    pic_w, pic_h = _fit_picture_box(px_w, px_h, avail_w, avail_h)
    left = (sw - pic_w) // 2
    slide.shapes.add_picture(
        stream, Emu(left), Emu(top), width=Emu(pic_w), height=Emu(pic_h)
    )


def add_title_slide(prs: Presentation, n_confirmed: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    box = slide.shapes.add_textbox(
        Emu(int(0.5 * EMU_PER_INCH)),
        Emu(sh // 2 - int(1.4 * EMU_PER_INCH)),
        Emu(sw - int(1.0 * EMU_PER_INCH)),
        Emu(int(2.8 * EMU_PER_INCH)),
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "FRB host GALFIT verification"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = (
        f"Population diagnostics  →  {n_confirmed} confirmed panels "
        f"(science cut: mag ≤ 22, b/a > 0.2)\n"
        f"Full-resolution PNGs embedded at {EMBED_DPI} DPI"
    )
    p2.font.size = Pt(15)
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.color.rgb = RGBColor(0x44, 0x44, 0x44)


def _fmt(v: float, nd: int = 3) -> str:
    try:
        return f"{float(v):.{nd}f}"
    except (TypeError, ValueError):
        return "—"


def add_population_overview_slide(prs: Presentation) -> None:
    """Text slide with key cohort null-test numbers (from population_summary.json)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    box = slide.shapes.add_textbox(
        Emu(int(0.45 * EMU_PER_INCH)),
        Emu(int(0.35 * EMU_PER_INCH)),
        Emu(sw - int(0.9 * EMU_PER_INCH)),
        Emu(sh - int(0.7 * EMU_PER_INCH)),
    )
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    lines: list[tuple[str, int, bool]] = [
        ("Population-level tests (in_53)", 22, True),
        ("Source: outputs/tables/population_summary.json + outputs/plots/", 11, False),
        ("", 10, False),
    ]
    if POP_SUMMARY.is_file():
        d = json.loads(POP_SUMMARY.read_text(encoding="utf-8"))
        lines.extend(
            [
                ("PSF leakage null tests", 14, True),
                (
                    f"  Spearman(q, e_PSF) = {_fmt(d.get('psf_q_vs_epsf_spearman_in53'), 3)}  "
                    f"(p = {_fmt(d.get('psf_q_vs_epsf_p_in53'), 3)})  — no amplitude correlation",
                    12,
                    False,
                ),
                (
                    f"  KS(|ΔPA|/90 vs Uniform) p = {_fmt(d.get('psf_dpa_ks_p_in53'), 3)}  "
                    f"(median |ΔPA| = {_fmt(d.get('psf_dpa_median_deg_in53'), 1)}°)  — flat PA alignment",
                    12,
                    False,
                ),
                ("", 8, False),
                ("Geometry agreement (production metrics)", 14, True),
                (
                    f"  sky Δq median [p16, p84] = {_fmt(d.get('dq_sky_median_in53'), 3)} "
                    f"[{_fmt(d.get('dq_sky_p16_in53'), 3)}, {_fmt(d.get('dq_sky_p84_in53'), 3)}]",
                    12,
                    False,
                ),
                (
                    f"  AstroPhot Δq median = {_fmt(d.get('dq_astrophot_median_in53'), 3)}  ·  "
                    f"iso Δq(2Re) median = {_fmt(d.get('iso_dq_2re_median_in53'), 3)}",
                    12,
                    False,
                ),
                (
                    f"  Fourier δq median = {_fmt(d.get('fourier_dq_median_in53'), 3)}  "
                    "(includes unreliable / unresolved hosts)",
                    12,
                    False,
                ),
                ("", 8, False),
                ("Photometry / cleanliness", 14, True),
                (
                    f"  Δm vs Re slope significance = {_fmt(d.get('dmag_vs_re_slope_sig_in53'), 1)}σ  "
                    f"(strong size–mag systematic)",
                    12,
                    False,
                ),
                (
                    f"  RFF(2Re) median = {_fmt(d.get('rff_2re_median_in53'), 3)}  ·  "
                    f"χ²/ν(2Re) median = {_fmt(d.get('chi2nu_local_2re_median_in53'), 2)}",
                    12,
                    False,
                ),
                ("", 8, False),
                (
                    "Next slides: one diagnostic figure each. Then confirmed host panels.",
                    12,
                    False,
                ),
            ]
        )
    else:
        lines.append(("population_summary.json missing — plots follow without numbers.", 12, False))

    first = True
    for text, size, bold in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = "Calibri"
        p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A) if bold else RGBColor(0x33, 0x33, 0x33)


def add_section_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    box = slide.shapes.add_textbox(
        Emu(int(0.5 * EMU_PER_INCH)),
        Emu(sh // 2 - int(0.8 * EMU_PER_INCH)),
        Emu(sw - int(1.0 * EMU_PER_INCH)),
        Emu(int(1.6 * EMU_PER_INCH)),
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.name = "Calibri"
        p2.alignment = PP_ALIGN.CENTER
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def main() -> int:
    hc = pd.read_csv(HC)
    hc["confirmed"] = hc["confirmed"].astype(str).str.lower().eq("true")
    cohort = vc.cohort("all64")[["frb", "mag", "b_a", "in_53"]]
    m = hc.merge(cohort, on="frb", how="left")
    m["in_53"] = m["in_53"].astype(bool)

    confirmed = m[m["confirmed"] & m["in_53"]].sort_values("frb")
    rejected = m[(~m["confirmed"]) & m["in_53"]].sort_values("frb")

    jobs: list[tuple[str, str, Path, str]] = []
    missing: list[str] = []
    for i, (_, row) in enumerate(confirmed.iterrows(), 1):
        frb = str(row["frb"])
        notes = "" if pd.isna(row.get("notes")) else str(row["notes"])
        img = panel_for(frb, notes)
        if not img.is_file():
            missing.append(frb)
            continue
        alt = _ALT.search(notes)
        leg = alt.group(1).split("_", 1)[1] if alt else "production"
        jobs.append(
            (f"{i}/{len(confirmed)}   {frb}   [{leg}]", "", img, "ok")
        )

    reject_jobs: list[tuple[str, str, Path]] = []
    for _, row in rejected.iterrows():
        frb = str(row["frb"])
        notes = "" if pd.isna(row.get("notes")) else str(row["notes"])
        compare = VER / "Re-fits" / frb / "psf_only" / "compare_sersic_vs_psf.png"
        if "star" in notes.lower() and compare.is_file():
            img = compare
        else:
            img = PANELS / f"{frb}.png"
        if img.is_file():
            reject_jobs.append(
                (
                    f"{frb}   REJECTED — not in confirmed sample",
                    (notes[:160] + "…") if len(notes) > 160 else notes,
                    img,
                )
            )

    pop_jobs: list[tuple[str, str, Path]] = []
    for fname, title, sub in POPULATION_SLIDES:
        path = PLOTS / fname
        if path.is_file():
            pop_jobs.append((title, sub, path))
        else:
            print(f"missing population plot: {path}", file=sys.stderr)

    probe = jobs[0][2] if jobs else (reject_jobs[0][2] if reject_jobs else None)
    if probe is None and not pop_jobs:
        print("no panels to embed", file=sys.stderr)
        return 1
    size_src = probe if probe is not None else pop_jobs[0][2]
    with Image.open(size_src) as im:
        px_w, px_h = im.size

    prs = Presentation()
    _configure_slide_size(prs, px_w, px_h)

    add_title_slide(prs, len(confirmed))
    add_population_overview_slide(prs)
    for title, sub, img in pop_jobs:
        add_image_slide(prs, title, img, subtitle=sub)

    add_section_slide(
        prs,
        "Confirmed host panels",
        f"{len(jobs)} in-cut confirmed  ·  then rejected in-cut hosts",
    )
    for title, _sub, img, _ in jobs:
        add_image_slide(prs, title, img)

    for title, sub, img in reject_jobs:
        add_image_slide(
            prs,
            title,
            img,
            title_color=RGBColor(0xA0, 0x20, 0x20),
            subtitle=sub,
        )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))

    verify = Presentation(str(OUT))
    sizes = [
        len(shape.image.blob)
        for slide in verify.slides
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    print(f"wrote {OUT}  ({OUT.stat().st_size // 1024} KB)")
    print(f"population slides: {1 + len(pop_jobs)}  (overview + figures)")
    print(f"confirmed slides: {len(jobs)}  missing: {missing}")
    print(f"rejected slides: {len(reject_jobs)}")
    print(f"total slides: {len(verify.slides)}")
    if sizes:
        print(
            f"embedded blobs: min={min(sizes) // 1024}KB  "
            f"med={sorted(sizes)[len(sizes) // 2] // 1024}KB  "
            f"max={max(sizes) // 1024}KB"
        )
    return 1 if missing else 0


if __name__ == "__main__":
    raise SystemExit(main())
